import os
import pandas as pd
from docx import Document
import PyPDF2
from pptx import Presentation
import openpyxl
import xlrd

def read_file(filepath):
    extension = os.path.splitext(filepath)[1].lower()

    if extension == '.txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()

    elif extension == '.docx':
        doc = Document(filepath)
        return '\n'.join([para.text for para in doc.paragraphs])

    elif extension == '.pdf':
        content = ""
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                content += page.extract_text() or ''
        return content

    elif extension == '.pptx':
        prs = Presentation(filepath)
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        return content
    
    elif extension in ['.xlsx', '.xls','.xlsm']:
            df = pd.read_excel(filepath, sheet_name=None)
            content = ""
            for sheet_name, sheet_data in df.items():
                content += f"\nSheet: {sheet_name}\n"
                content += sheet_data.to_string(index=False) + "\n"
            return content

    elif extension == '.csv':
            df = pd.read_csv(filepath)
            return df.to_string(index=False)
    else:
        return f"Unsupported file type: {extension}"


content=read_file(r"S:\cbe\VPA_ASSIGMENT\vba\test.xlsm")
print(content)